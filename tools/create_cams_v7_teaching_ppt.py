from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = Path.cwd() / "CAMS_V7_教研沟通版_流程图.pptx"


class C:
    ink = "1F2933"
    muted = "64748B"
    slate = "334155"
    line = "CBD5E1"
    pale = "F7FAFC"
    panel = "FFFFFF"
    green = "0F766E"
    green2 = "14B8A6"
    teal_pale = "D9F3F0"
    amber = "F59E0B"
    amber_pale = "FEF3C7"
    blue = "2563EB"
    blue_pale = "DBEAFE"
    purple = "7C3AED"
    purple_pale = "EDE9FE"
    dark = "102A43"
    dark2 = "16324F"
    white = "FFFFFF"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def set_text(
    shape,
    text: str,
    size: float = 12,
    color: str = C.ink,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    font: str = "Microsoft YaHei",
):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return shape


def add_text(slide, text, x, y, w, h, **kwargs):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    return set_text(shape, text, **kwargs)


def add_rect(slide, x, y, w, h, fill=C.panel, line="E2E8F0", radius=False):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.line.color.rgb = rgb(line)
    shp.line.width = Pt(0.8)
    return shp


def add_card(slide, x, y, w, h, fill=C.panel, line="E2E8F0"):
    return add_rect(slide, x, y, w, h, fill=fill, line=line, radius=True)


def add_title(slide, title, sub=None):
    add_text(slide, title, 0.55, 0.35, 8.8, 0.45, size=24, color=C.ink, bold=True)
    if sub:
        add_text(slide, sub, 0.56, 0.86, 10.8, 0.28, size=10.5, color=C.muted)


def add_footer(slide, page):
    add_text(slide, "依据：技术路线总图.md / CAMSV7 260630 材料", 0.55, 7.12, 9.5, 0.18, size=8.5, color="94A3B8")
    add_text(slide, f"{page:02d}", 12.15, 7.03, 0.55, 0.18, size=9, color="94A3B8", bold=True, align=PP_ALIGN.RIGHT)


def add_pill(slide, text, x, y, w, color, fill="FFFFFF"):
    add_card(slide, x, y, w, 0.34, fill=fill, line=color)
    add_text(slide, text, x + 0.05, y + 0.07, w - 0.1, 0.13, size=8.8, color=color, bold=True, align=PP_ALIGN.CENTER)


def add_arrow(slide, x1, y1, x2, y2, color=C.line):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(1.3)
    line.line.end_arrowhead = True
    return line


def add_step(slide, x, y, w, h, num, header, body, color, fill):
    add_card(slide, x, y, w, h, fill=fill, line=color)
    circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.18), Inches(y + 0.18), Inches(0.42), Inches(0.42))
    circ.fill.solid()
    circ.fill.fore_color.rgb = rgb("FFFFFF")
    circ.line.color.rgb = rgb("FFFFFF")
    add_text(slide, str(num), x + 0.18, y + 0.285, 0.42, 0.12, size=9, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, header, x + 0.72, y + 0.17, w - 0.85, 0.24, size=10.6, color=color, bold=True)
    add_text(slide, body, x + 0.72, y + 0.48, w - 0.85, h - 0.55, size=8.4, color=C.slate)


def add_info_box(slide, x, y, w, h, header, body, color, fill):
    add_card(slide, x, y, w, h, fill=fill, line=color)
    add_text(slide, header, x + 0.2, y + 0.14, w - 0.4, 0.26, size=12.4, color=color, bold=True)
    add_text(slide, body, x + 0.2, y + 0.51, w - 0.4, h - 0.62, size=9.2, color=C.slate)


def set_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(color)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# Slide 1
slide = prs.slides.add_slide(blank)
set_bg(slide, C.dark)
add_rect(slide, 8.62, 0, 4.70, 7.5, fill=C.green, line=C.green)
add_text(slide, "CAMS V7\n教研工作台沟通版", 0.8, 1.05, 6.85, 1.55, size=33, color=C.white, bold=True)
add_text(slide, "以教材为底座，把官方题、解析、答疑都挂回知识点", 0.84, 2.85, 7.05, 0.35, size=16, color=C.teal_pale)
add_text(slide, "明天沟通重点：先讲全流程与最终产物，工程细节一笔带过。", 0.86, 3.45, 6.8, 0.42, size=13, color="B6E6DF")
for i, (a, b) in enumerate([("教材", "句卡"), ("V7题", "绑定"), ("解析", "复核"), ("答疑", "沉淀")]):
    y = 1.1 + i * 1.08
    add_card(slide, 9.15, y, 1.3, 0.56, fill="FFFFFF", line="FFFFFF")
    add_text(slide, a, 9.15, y + 0.16, 1.3, 0.14, size=10, color=C.green, bold=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, 10.55, y + 0.28, 11.18, y + 0.28, "B6E6DF")
    add_card(slide, 11.35, y, 1.3, 0.56, fill=C.amber_pale, line=C.amber_pale)
    add_text(slide, b, 11.35, y + 0.16, 1.3, 0.14, size=10, color="92400E", bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "260630", 0.86, 6.6, 1.3, 0.18, size=10, color="94A3B8")

# Slide 2
slide = prs.slides.add_slide(blank)
set_bg(slide, C.pale)
add_title(slide, "先把口径对齐", "教研需要先理解“系统如何产生价值”，不是先听脚本和字段。")
add_info_box(slide, 0.65, 1.55, 3.75, 3.7, "一句话定位", "这不是单纯让 AI 写解析，而是以教材为中心，把 V7 官方题、解析、学生答疑都关联回教材知识点的教研工作台。", C.green, C.teal_pale)
add_info_box(slide, 4.8, 1.55, 3.75, 3.7, "第一阶段目标", "先服务 V7 官方题库：题目整理、证据绑定、候选解析、教研复核、上架。标准是短、准、能回到知识点。", C.blue, C.blue_pale)
add_info_box(slide, 8.95, 1.55, 3.75, 3.7, "AI 的边界", "AI 负责证据寻找和候选答案/解析生成；教研负责判断证据是否合适、解析是否准确、内容是否值得沉淀。", C.amber, C.amber_pale)
add_pill(slide, "明天讲法：先全流程，再看最终产物，最后确认教研拍板点", 3.0, 6.05, 7.35, C.green)
add_footer(slide, 2)

# Slide 3
slide = prs.slides.add_slide(blank)
set_bg(slide, "FFFFFF")
add_title(slide, "全流程总图", "三条输入流都回到同一个教材证据底座：教材句卡与知识点。")
y1, y2, y3 = 1.35, 3.0, 4.65
add_step(slide, 0.55, y1, 2.1, 0.92, "1", "教材 PDF", "MinerU + 清洗\n结构化 Markdown", C.green, C.teal_pale)
add_step(slide, 3.05, y1, 2.1, 0.92, "2", "教材句卡", "5199 张\n可追溯原文证据", C.green, "EFFCFB")
add_step(slide, 5.55, y1, 2.1, 0.92, "3", "教材知识图谱", "知识点 / 法规 / 风险\n关系与句卡挂载", C.green, C.teal_pale)
add_step(slide, 0.55, y2, 2.1, 0.92, "A", "V7 官方题库", "题干 / 选项 / 答案\n结构化入库", C.blue, C.blue_pale)
add_step(slide, 3.05, y2, 2.1, 0.92, "B", "题目-选项-句卡绑定", "检索教材证据\n输出 q_{id}.json", C.blue, "EFF6FF")
add_step(slide, 5.55, y2, 2.1, 0.92, "C", "AI 候选解析", "候选答案\n正确/错误项解释", C.blue, C.blue_pale)
add_step(slide, 8.05, y2, 2.1, 0.92, "D", "教研复核", "证据、解析、易错点\n人工定稿", C.amber, C.amber_pale)
add_step(slide, 10.55, y2, 2.1, 0.92, "E", "上架与反馈", "解析上线\n用户反馈再迭代", C.amber, "FFF7ED")
add_step(slide, 0.55, y3, 2.1, 0.92, "α", "学生答疑", "问题 / 对应题目\n结构化记录", C.purple, C.purple_pale)
add_step(slide, 3.05, y3, 2.1, 0.92, "β", "候选答疑", "基于题目与教材证据\n生成初稿", C.purple, "F5F3FF")
add_step(slide, 5.55, y3, 2.1, 0.92, "γ", "是否沉淀", "有代表性才回挂\n避免知识库变乱", C.purple, C.purple_pale)
add_step(slide, 8.05, y3, 2.1, 0.92, "δ", "反哺教研", "易错点 / FAQ\n解析优化", C.purple, "F5F3FF")
for a in [
    (2.65, y1 + 0.46, 3.04, y1 + 0.46), (5.15, y1 + 0.46, 5.54, y1 + 0.46),
    (2.65, y2 + 0.46, 3.04, y2 + 0.46), (5.15, y2 + 0.46, 5.54, y2 + 0.46),
    (7.65, y2 + 0.46, 8.04, y2 + 0.46), (10.15, y2 + 0.46, 10.54, y2 + 0.46),
    (2.65, y3 + 0.46, 3.04, y3 + 0.46), (5.15, y3 + 0.46, 5.54, y3 + 0.46),
    (7.65, y3 + 0.46, 8.04, y3 + 0.46),
]:
    add_arrow(slide, *a, C.line)
add_arrow(slide, 6.55, 2.25, 4.35, 2.95, C.green2)
add_arrow(slide, 4.1, 3.94, 4.1, 4.62, C.purple)
add_text(slide, "工程细节可以一句话带过：底层已把教材拆成可追溯句卡，并用检索 + 裁判生成每题绑定过程。", 0.65, 6.22, 11.6, 0.26, size=10.8, color=C.muted)
add_footer(slide, 3)

# Slide 4
slide = prs.slides.add_slide(blank)
set_bg(slide, C.pale)
add_title(slide, "最终产物：一题一份可复核解析包", "技术主产物是 q_{id}.json；教研看到的应是可读的题目解析包。")
add_card(slide, 0.7, 1.25, 5.25, 5.15, fill="FFFFFF", line="CBD5E1")
add_text(slide, "示例：第 N 题解析包", 1.05, 1.55, 3.8, 0.25, size=17, color=C.ink, bold=True)
add_pill(slide, "主产物：output/questions/q_{id}.json", 1.05, 2.0, 3.1, C.green, C.teal_pale)
rows = [
    ("题目与答案", "题干、选项、标准答案、题型"),
    ("教材证据", "每个选项对应的强证据 / 候选句卡"),
    ("AI 候选解析", "候选答案、正确项说明、错误项解释"),
    ("知识点位置", "教材章节、句卡、知识图谱节点"),
    ("复核标记", "证据不足、分歧、多 direct 冲突"),
]
for i, (h, b) in enumerate(rows):
    y = 2.55 + i * 0.62
    add_rect(slide, 1.05, y, 0.12, 0.36, fill=C.blue if i % 2 else C.green, line=C.blue if i % 2 else C.green)
    add_text(slide, h, 1.3, y + 0.04, 1.3, 0.16, size=10.8, color=C.ink, bold=True)
    add_text(slide, b, 2.65, y + 0.04, 2.85, 0.16, size=9.6, color=C.slate)
add_info_box(slide, 6.45, 1.25, 2.85, 1.55, "解析标准", "短、准、能回到知识点；不写教材扩写，也不只写“答案是 A”。", C.green, C.teal_pale)
add_info_box(slide, 9.65, 1.25, 2.85, 1.55, "证据标准", "每个关键判断尽量能回到教材句卡；弱证据和缺证据要显式标记。", C.blue, C.blue_pale)
add_info_box(slide, 6.45, 3.1, 2.85, 1.55, "复核标准", "AI 候选答案只是候选；教研确认后才进入正式解析。", C.amber, C.amber_pale)
add_info_box(slide, 9.65, 3.1, 2.85, 1.55, "沉淀标准", "有代表性的错误项、答疑、易错点沉淀回知识点。", C.purple, C.purple_pale)
add_text(slide, "教研最终不需要看 JSON 字段，而是看一份能审核、能改、能上架的题目解析包。", 6.45, 5.42, 5.8, 0.38, size=14, color=C.ink, bold=True)
add_footer(slide, 4)

# Slide 5
slide = prs.slides.add_slide(blank)
set_bg(slide, "FFFFFF")
add_title(slide, "教研需要介入的地方", "系统越自动，越要把人工判断入口设计清楚。")
add_text(slide, "AI 先做", 0.8, 1.25, 2.0, 0.28, size=18, color=C.green, bold=True)
add_text(slide, "教研拍板", 7.0, 1.25, 2.0, 0.28, size=18, color=C.amber, bold=True)
for i, (h, b) in enumerate([
    ("找教材证据", "从句卡和知识图谱中召回候选依据"),
    ("生成候选答案", "基于证据判断选项对错"),
    ("写候选解析", "解释正确项、错误项和疑似误区"),
    ("标记风险", "缺证据、答案分歧、证据冲突"),
]):
    add_step(slide, 0.8, 1.75 + i * 1.05, 4.65, 0.72, i + 1, h, b, C.green, C.teal_pale)
for i, (h, b) in enumerate([
    ("证据是否合适", "这张句卡能不能支撑这个选项"),
    ("解析是否准确", "是否短、准、讲到学生会错的点"),
    ("知识点怎么命名", "章节名、知识点名、考点名的口径"),
    ("是否值得沉淀", "答疑/易错点是否回挂教材"),
]):
    add_step(slide, 7.0, 1.75 + i * 1.05, 4.65, 0.72, i + 1, h, b, C.amber, C.amber_pale)
add_arrow(slide, 5.62, 3.63, 6.82, 3.63, C.line)
add_text(slide, "核心原则：AI 给候选，教研定结论。", 3.55, 6.25, 6.1, 0.35, size=17, color=C.ink, bold=True, align=PP_ALIGN.CENTER)
add_footer(slide, 5)

# Slide 6
slide = prs.slides.add_slide(blank)
set_bg(slide, C.pale)
add_title(slide, "明天建议直接确认的 5 件事", "把沟通从“系统介绍”推进到“样板怎么跑、怎么验收”。")
for i, (h, b) in enumerate([
    ("解析模板", "是否固定为：考什么 / 为什么正确 / 为什么错误 / 对应知识点 / 易错点"),
    ("长度标准", "每题解析控制在多少字；单选、多选、判断题是否不同"),
    ("知识点口径", "用教材章节、KG 节点，还是教研确认后的考点名"),
    ("复核方式", "全量逐题审核、低置信优先审核，还是先抽样校准"),
    ("样板范围", "先选 20-50 道 V7 官方题跑通样板，再批量展开"),
]):
    x = 0.75 + i * 4.15 if i < 3 else 2.9 + (i - 3) * 4.15
    y = 1.45 if i < 3 else 4.0
    add_info_box(slide, x, y, 3.55, 1.45, f"{i + 1}. {h}", b, C.blue if i % 2 else C.green, C.blue_pale if i % 2 else C.teal_pale)
add_text(slide, "建议收口：先出一批可审核样板，让教研校准模板和证据标准，再批量处理 V7 官方题。", 1.35, 6.28, 10.7, 0.32, size=14, color=C.ink, bold=True, align=PP_ALIGN.CENTER)
add_footer(slide, 6)

# Slide 7
slide = prs.slides.add_slide(blank)
set_bg(slide, C.dark2)
add_text(slide, "如果教研追问工程细节", 0.75, 0.75, 5.5, 0.45, size=25, color=C.white, bold=True)
add_text(slide, "一笔带过即可，不要让沟通陷进字段和脚本。", 0.78, 1.28, 5.9, 0.28, size=12.5, color="B6E6DF")
for i, (h, b) in enumerate([
    ("教材底座", "PDF -> 结构化 Markdown -> 5199 张句卡，句卡可回到教材原文和页码。"),
    ("知识图谱", "已完成教材 KG：节点、关系、句卡挂载，用于导航和辅助召回。"),
    ("题目绑定", "V7 题进来后，主产物是一题一份 q_{id}.json，记录证据寻找全过程。"),
    ("考点沉淀", "考点生成仍是 preview，等题目证据边稳定后再做归并和命名。"),
]):
    y = 2.0 + i * 1.05
    add_card(slide, 0.82, y, 11.7, 0.72, fill="1E3A5F" if i % 2 else "173B57", line="2B5876")
    add_text(slide, h, 1.18, y + 0.21, 1.55, 0.18, size=11.5, color="FDE68A", bold=True)
    add_text(slide, b, 2.85, y + 0.21, 8.8, 0.2, size=10.8, color="E2E8F0")
add_text(slide, "推荐回应：底层能力已经在支撑“短、准、能回到知识点”的解析生产；明天先确认教研验收标准。", 0.82, 6.45, 11.2, 0.3, size=13, color=C.white, bold=True)

prs.save(OUT)
print(OUT)
