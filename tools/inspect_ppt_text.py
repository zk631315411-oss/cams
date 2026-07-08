from pathlib import Path

from pptx import Presentation


p = next(Path.cwd().glob("CAMS_V7*.pptx"))
prs = Presentation(p)
print("file", p)
print("slides", len(prs.slides))
slide_w = prs.slide_width
slide_h = prs.slide_height
issues = []
for i, slide in enumerate(prs.slides, 1):
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            texts.append(shape.text.strip().replace("\n", " / "))
        if shape.left < 0 or shape.top < 0 or shape.left + shape.width > slide_w or shape.top + shape.height > slide_h:
            issues.append(
                f"slide {i}: shape out of bounds "
                f"type={shape.shape_type} left={shape.left} top={shape.top} width={shape.width} height={shape.height}"
            )
    print(f"--- {i} ---")
    print(" | ".join(texts[:20]))
print("geometry_issues", len(issues))
for issue in issues[:20]:
    print(issue)
